import os
import re

from colorama import Fore

from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
from lxml import etree
from bs4 import BeautifulSoup
from markdownify import MarkdownConverter

from ollama_chat.core import plugins

def extract_text_from_pdf(pdf_content):
    with open('temp.pdf', 'wb') as f:
        f.write(pdf_content)

    reader = PdfReader('temp.pdf')
    text = ''
    for page in reader.pages:
        text += page.extract_text()

    # Clean up by removing the temporary file
    os.remove('temp.pdf')

    # Return the extracted text, with extra newlines removed
    return re.sub(r'\n+', '\n', text)

def extract_text_from_docx(docx_path):
    # Load the Word document
    document = Document(docx_path)

    # Extract the file name (without extension) and replace underscores with spaces
    file_name = os.path.splitext(os.path.basename(docx_path))[0].replace('_', ' ')

    # Initialize a list to collect Markdown lines
    markdown_lines = []

    def process_paragraph(paragraph, list_level=0):
        """Convert a paragraph into Markdown based on its style and list level."""
        text = paragraph.text.replace("\n", " ").strip()  # Replace carriage returns with spaces
        if not text:
            return None  # Skip empty paragraphs

        # Check if paragraph is a list item based on indentation
        if paragraph.style.name == "List Paragraph":
            # Use the list level to determine indentation for bullet points
            bullet_prefix = "  " * list_level + "- "
            return f"{bullet_prefix}{text}"

        # Check for headings
        if paragraph.style.name.startswith("Heading"):
            heading_level = int(paragraph.style.name.split(" ")[1])
            return f"{'#' * heading_level} {text}"

        # Default: Regular paragraph
        return text

    def extract_lists(docx):
        """Extract the list structure from the underlying XML of the document."""
        # Access the document XML using lxml
        xml_content = docx.element
        namespaces = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        }

        # Parse the XML tree using lxml's etree
        root = etree.fromstring(etree.tostring(xml_content))

        # Find all list items (w:li)
        list_paragraphs = []
        for item in root.xpath("//w:li", namespaces=namespaces):
            # Extract the list level from the parent elements
            list_level = item.getparent().getparent().get("w:ilvl")
            if list_level is not None:
                list_paragraphs.append((list_level, item.text.strip()))

        return list_paragraphs

    # Add the document title (file name) as the top-level heading
    markdown_lines.append(f"# {file_name}")

    # Process each paragraph in the document
    for paragraph in document.paragraphs:
        # Detect bullet points based on paragraph's indent level (style `List Paragraph`)
        markdown_line = process_paragraph(paragraph)
        if markdown_line:
            markdown_lines.append(markdown_line)

    # Extract and process lists directly from the document's XML
    lists = extract_lists(document)
    for level, item in lists:
        bullet_prefix = "  " * int(level) + "- "
        markdown_lines.append(f"{bullet_prefix}{item}")

    # Join all lines into a single Markdown string
    return "\n\n".join(markdown_lines)

def extract_text_from_pptx(pptx_path):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_path)

    # Extract the file name (without extension) and replace underscores with spaces
    file_name = os.path.splitext(os.path.basename(pptx_path))[0].replace('_', ' ')

    # Initialize a list to collect Markdown lines
    markdown_lines = []

    def extract_text_with_bullets(shape, exclude_text=None):
        """Extract text with proper bullet point levels from a shape."""
        text_lines = []
        if shape.is_placeholder or shape.has_text_frame:
            if shape.text_frame and shape.text_frame.text.strip():
                for paragraph in shape.text_frame.paragraphs:
                    line_text = paragraph.text.replace("\r", "").replace("\n", " ").strip()  # Replace \n with space
                    if line_text and line_text != exclude_text:  # Exclude the slide title if needed
                        bullet_level = paragraph.level  # Get the bullet level
                        bullet = "  " * bullet_level + "- " + line_text
                        text_lines.append(bullet)
        elif shape.shape_type == 6:  # Grouped shapes
            # Handle grouped shapes recursively
            for sub_shape in shape.shapes:
                text_lines.extend(extract_text_with_bullets(sub_shape, exclude_text))
        return text_lines

    def get_first_text_entry(slide):
        """Retrieve the first text entry from the slide."""
        for shape in slide.shapes:
            if shape.is_placeholder or shape.has_text_frame:
                if shape.text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.paragraphs[0].text.replace("\n", " ").strip()
        return None

    for slide_number, slide in enumerate(presentation.slides, start=1):
        # Determine the Markdown header level
        if slide_number == 1:
            header_prefix = "#"
        else:
            header_prefix = "##"

        # Add the slide title or file name as the main title for the first slide
        if slide_number == 1:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
            else:
                title = file_name
            markdown_lines.append(f"{header_prefix} {title}")
        else:
            # Add the title for subsequent slides
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
            else:
                # Use the first text entry as the slide title if no title is present
                title = get_first_text_entry(slide)
                if not title:
                    title = f"Slide {slide_number}"
            markdown_lines.append(f"{header_prefix} {title}")

        # Add the slide content (text in other shapes), excluding the title if it's used
        for shape in slide.shapes:
            bullet_text = extract_text_with_bullets(shape, exclude_text=title)
            markdown_lines.extend(bullet_text)

        # Add a separator between slides, except after the last slide
        if slide_number < len(presentation.slides):
            markdown_lines.append("")

    # Join all lines into a single Markdown string
    return "\n".join(markdown_lines)

def extract_text_from_html(html_content):
    # Convert the modified HTML content to Markdown
    try:
        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove all <script> tags
        for script in soup.find_all('script'):
            script.decompose()

        # Remove all <style> tags
        for style in soup.find_all('style'):
            style.decompose()

        # Remove all <noscript> tags
        for noscript in soup.find_all('noscript'):
            noscript.decompose()

        # Remove all <svg> tags
        for svg in soup.find_all('svg'):
            svg.decompose()

        # Remove all <canvas> tags
        for canvas in soup.find_all('canvas'):
            canvas.decompose()

        # Remove all <audio> tags
        for audio in soup.find_all('audio'):
            audio.decompose()

        # Remove all <video> tags
        for video in soup.find_all('video'):
            video.decompose()

        # Remove all <iframe> tags
        for iframe in soup.find_all('iframe'):
            iframe.decompose()

        text = md(soup, strip=['a', 'img'], heading_style='ATX',
                        escape_asterisks=False, escape_underscores=False,
                        autolinks=False)

        # Remove extra newlines
        text = re.sub(r'\n+', '\n', text)

        return text
    except Exception as e:
        plugins.on_print(f"Failed to parse HTML content: {e}", Fore.RED)
        return ""

def md(soup, **options):
    return MarkdownConverter(**options).convert_soup(soup)
